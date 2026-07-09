from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paleta de cores ──────────────────────────────────────────────
PRETO       = RGBColor(0x0D, 0x0D, 0x0D)
CINZA_ESC   = RGBColor(0x1A, 0x1A, 0x2E)
CINZA_MED   = RGBColor(0x16, 0x21, 0x3E)
AZUL_ESC    = RGBColor(0x0F, 0x3C, 0x60)
VERDE       = RGBColor(0x00, 0xD4, 0x8A)
VERDE_ESC   = RGBColor(0x00, 0xA8, 0x6B)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TEXTO = RGBColor(0xB0, 0xB8, 0xC8)
AMARELO     = RGBColor(0xFF, 0xD7, 0x00)
LARANJA     = RGBColor(0xFF, 0x6B, 0x35)
AZUL_CARD   = RGBColor(0x0D, 0x2A, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=24, bold=False, color=BRANCO,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_text(slide, lines, l, t, w, h, size=14, color=BRANCO, bold_first=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and idx == 0)
    return txBox


def bg(slide):
    add_rect(slide, 0, 0, W, H, CINZA_ESC)
    add_rect(slide, 0, 0, W, Inches(0.08), VERDE)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), VERDE)


def accent_bar(slide):
    add_rect(slide, 0, 0, Inches(0.12), H, VERDE)


def slide_header(slide, icon_title, subtitle=None):
    add_text(slide, icon_title, Inches(0.5), Inches(0.35), Inches(12), Inches(0.72),
             size=34, bold=True, color=VERDE)
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.5), Inches(0.045), VERDE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
                 size=15, color=CINZA_TEXTO)


def numero_slide(slide, num, total=11):
    add_text(slide, f"{num} / {total}", Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.35),
             size=12, color=CINZA_TEXTO, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)

add_rect(s, 0, Inches(2.1), W, Inches(3.3), CINZA_MED)
add_rect(s, 0, Inches(2.1), Inches(0.18), Inches(3.3), VERDE)

add_text(s, "🏋️", Inches(0.5), Inches(0.5), Inches(2.5), Inches(1.5), size=80, align=PP_ALIGN.CENTER)

add_text(s, "FIT TRACK", Inches(2.2), Inches(2.22), Inches(9.0), Inches(1.3),
         size=72, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

add_text(s, "Treine. Evolua. Supere.", Inches(2.2), Inches(3.38), Inches(9.0), Inches(0.7),
         size=24, italic=True, color=CINZA_TEXTO, align=PP_ALIGN.CENTER)

add_text(s, "Aplicativo Mobile de Acompanhamento Fitness",
         Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6),
         size=19, color=BRANCO, align=PP_ALIGN.CENTER)

add_rect(s, Inches(4.2), Inches(4.95), Inches(4.93), Inches(0.045), VERDE)

add_text(s, "4º Semestre — Análise e Desenvolvimento de Sistemas  |  2026",
         Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.5),
         size=15, color=CINZA_TEXTO, align=PP_ALIGN.CENTER)

numero_slide(s, 1)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Integrantes
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "👥  Integrantes do Grupo")
numero_slide(s, 2)

members = [
    ("Lucas Felipe",    "Frontend & UI Design"),
    ("Kaylan Luis",     "Backend & Firebase"),
    ("Matheus Antony",  "Gamificação & Missões"),
    ("Luis Henrique",   "Navegação & Estado Global"),
]

for i, (name, role) in enumerate(members):
    col = i % 2
    row = i // 2
    lft = Inches(0.9 + col * 6.2)
    top = Inches(1.7 + row * 2.3)

    add_rect(s, lft, top, Inches(5.7), Inches(1.95), AZUL_CARD)
    add_rect(s, lft, top, Inches(0.14), Inches(1.95), VERDE)

    # Número / avatar simulado
    add_rect(s, lft + Inches(0.3), top + Inches(0.4), Inches(0.85), Inches(0.85),
             RGBColor(0x00, 0xA8, 0x6B))
    add_text(s, str(i + 1), lft + Inches(0.3), top + Inches(0.38), Inches(0.85), Inches(0.85),
             size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    add_text(s, name, lft + Inches(1.35), top + Inches(0.3), Inches(4.0), Inches(0.65),
             size=22, bold=True, color=BRANCO)
    add_text(s, role, lft + Inches(1.35), top + Inches(0.95), Inches(4.0), Inches(0.5),
             size=15, color=VERDE)

add_text(s, "Curso: Análise e Desenvolvimento de Sistemas  |  4º Semestre",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
         size=13, color=CINZA_TEXTO, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Divisão de Tarefas
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "📋  Divisão de Tarefas",
             "Cada integrante foi responsável por uma área específica do projeto")
numero_slide(s, 3)

task_data = [
    ("Lucas Felipe",   "Frontend & UI",
     ["• Telas: Splash, Onboarding, Login/Cadastro",
      "• Dashboard, Layout geral e componentes",
      "• Tela de Execução de Treino em tempo real",
      "• Design system, cores e tipografia"]),
    ("Kaylan Luis",    "Backend & Firebase",
     ["• Configuração do Firebase (Auth, Firestore, Storage)",
      "• Serviços de autenticação (login, cadastro, logout)",
      "• CRUD de treinos e histórico no Firestore",
      "• Upload e gerenciamento de foto de perfil"]),
    ("Matheus Antony", "Gamificação & Missões",
     ["• Sistema de XP e progressão de níveis",
      "• Missões diárias e reset automático",
      "• 6 Badges / conquistas com critérios",
      "• Tela de Estatísticas e gráfico semanal"]),
    ("Luis Henrique",  "Navegação & Estado",
     ["• React Navigation: Stack + Bottom Tabs",
      "• Context API (AppContext) — estado global",
      "• AsyncStorage para persistência offline",
      "• Tela de Perfil, Histórico e Configurações"]),
]

CORES_CARD = [
    RGBColor(0x00, 0xA8, 0x6B),
    RGBColor(0x0F, 0x5C, 0x99),
    RGBColor(0xB0, 0x5E, 0x00),
    RGBColor(0x7B, 0x2D, 0x8B),
]

for i, (name, area, tasks) in enumerate(task_data):
    col = i % 2
    row = i // 2
    lft = Inches(0.5 + col * 6.4)
    top = Inches(1.65 + row * 2.75)

    add_rect(s, lft, top, Inches(5.9), Inches(2.5), AZUL_CARD)
    add_rect(s, lft, top, Inches(0.14), Inches(2.5), CORES_CARD[i])

    add_text(s, name, lft + Inches(0.28), top + Inches(0.1), Inches(5.3), Inches(0.5),
             size=17, bold=True, color=BRANCO)
    add_text(s, area, lft + Inches(0.28), top + Inches(0.55), Inches(5.3), Inches(0.38),
             size=13, color=VERDE, bold=True)
    add_rect(s, lft + Inches(0.28), top + Inches(0.93), Inches(5.4), Inches(0.03), CINZA_TEXTO)

    for j, task in enumerate(tasks):
        add_text(s, task, lft + Inches(0.28), top + Inches(0.98 + j * 0.36),
                 Inches(5.4), Inches(0.34), size=12, color=BRANCO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — O que é o FIT TRACK?
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "💡  O que é o FIT TRACK?",
             "Um aplicativo mobile completo para acompanhamento de treinos e evolução física pessoal.")
numero_slide(s, 4)

cards = [
    ("🎯", "Objetivo",       "Ajudar usuários a manter\nconsistência nos treinos"),
    ("📱", "Plataforma",     "Mobile — Android\n(React Native + Expo)"),
    ("☁️", "Backend",       "Firebase (Auth, Firestore,\nStorage)"),
    ("🎮", "Diferencial",   "Gamificação com XP,\nníveis, missões e conquistas"),
]

for i, (icon, title, desc) in enumerate(cards):
    col = i % 2
    row = i // 2
    lft = Inches(0.5 + col * 6.4)
    top = Inches(2.1 + row * 2.35)

    add_rect(s, lft, top, Inches(5.9), Inches(2.05), AZUL_CARD)
    add_rect(s, lft, top, Inches(0.14), Inches(2.05), VERDE)
    add_text(s, icon, lft + Inches(0.28), top + Inches(0.15), Inches(0.9), Inches(0.9), size=34)
    add_text(s, title, lft + Inches(1.2), top + Inches(0.12), Inches(4.4), Inches(0.55),
             size=19, bold=True, color=VERDE)
    add_text(s, desc, lft + Inches(1.2), top + Inches(0.7), Inches(4.4), Inches(1.1),
             size=16, color=BRANCO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Tecnologias
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "🛠️  Tecnologias Utilizadas")
numero_slide(s, 5)

techs = [
    ("⚛️",  "React Native 0.71",   "Framework mobile cross-platform"),
    ("📦",  "Expo ~48",             "Build & desenvolvimento ágil"),
    ("🔥",  "Firebase 12",          "Auth + Firestore + Storage"),
    ("🧭",  "React Navigation v6",  "Roteamento de telas (Stack + Tabs)"),
    ("💾",  "AsyncStorage",         "Persistência local / offline"),
    ("🖼️", "Expo Image Picker",    "Upload de foto de perfil"),
    ("🌐",  "Context API",          "Gerenciamento de estado global"),
    ("🎨",  "Theme System",         "Design system centralizado"),
]

for i, (icon, name, desc) in enumerate(techs):
    col = i % 2
    row = i // 2
    lft = Inches(0.5 + col * 6.4)
    top = Inches(1.3 + row * 1.45)

    add_rect(s, lft, top, Inches(5.9), Inches(1.25), AZUL_CARD)
    add_rect(s, lft, top, Inches(0.12), Inches(1.25), VERDE)
    add_text(s, icon, lft + Inches(0.22), top + Inches(0.16), Inches(0.7), Inches(0.7), size=24)
    add_text(s, name, lft + Inches(1.05), top + Inches(0.1), Inches(4.5), Inches(0.52),
             size=17, bold=True, color=VERDE)
    add_text(s, desc, lft + Inches(1.05), top + Inches(0.6), Inches(4.5), Inches(0.55),
             size=14, color=CINZA_TEXTO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Telas do Aplicativo
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "📱  Telas do Aplicativo", "11 telas cobrindo o fluxo completo de uso do aplicativo")
numero_slide(s, 6)

screens = [
    ("🚀", "Splash / Onboarding", "Inicialização e boas-vindas"),
    ("🔐", "Login / Cadastro",    "Auth com e-mail e senha"),
    ("🏠", "Dashboard",           "Visão geral, streak, missões"),
    ("💪", "Treinos",             "Biblioteca + criação customizada"),
    ("▶️", "Execução de Treino",  "Progressão exercício a exercício"),
    ("🎯", "Missões",             "XP, nível e desafios diários"),
    ("📊", "Estatísticas",        "Gráficos semanais e conquistas"),
    ("👤", "Perfil & Config.",    "Foto, dados pessoais, preferências"),
    ("📋", "Histórico",           "Log cronológico de treinos"),
]

for i, (icon, name, desc) in enumerate(screens):
    col = i % 3
    row = i // 3
    lft = Inches(0.5 + col * 4.25)
    top = Inches(1.5 + row * 1.9)

    add_rect(s, lft, top, Inches(3.9), Inches(1.65), AZUL_CARD)
    add_rect(s, lft, top, Inches(0.12), Inches(1.65), VERDE)
    add_text(s, icon + "  " + name, lft + Inches(0.24), top + Inches(0.14),
             Inches(3.5), Inches(0.58), size=15, bold=True, color=BRANCO)
    add_rect(s, lft + Inches(0.24), top + Inches(0.74), Inches(3.45), Inches(0.03), CINZA_TEXTO)
    add_text(s, desc, lft + Inches(0.24), top + Inches(0.82),
             Inches(3.5), Inches(0.7), size=13, color=CINZA_TEXTO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Funcionalidades Principais
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "⚡  Funcionalidades Principais")
numero_slide(s, 7)

funcs = [
    ("🔐", "Autenticação Firebase",
     "Cadastro, login e persistência automática de sessão com e-mail e senha."),
    ("💪", "Gerenciamento de Treinos",
     "5 treinos pré-configurados (Superior, Inferior, Push, Pull, HIIT) + criação personalizada."),
    ("▶️", "Execução em Tempo Real",
     "Progressão por exercício, contagem de séries e barra de progresso visual."),
    ("📸", "Upload de Foto de Perfil",
     "Câmera ou galeria via Expo Image Picker, armazenamento no Firebase Storage."),
]

for i, (icon, title, desc) in enumerate(funcs):
    top = Inches(1.35 + i * 1.48)

    add_rect(s, Inches(0.5), top, Inches(12.3), Inches(1.28), AZUL_CARD)
    add_rect(s, Inches(0.5), top, Inches(0.14), Inches(1.28), VERDE)

    # Ícone + título na esquerda
    add_text(s, icon, Inches(0.75), top + Inches(0.2), Inches(0.75), Inches(0.75), size=28)
    add_text(s, title, Inches(1.6), top + Inches(0.12), Inches(4.2), Inches(0.52),
             size=18, bold=True, color=VERDE)

    # Separador vertical
    add_rect(s, Inches(6.1), top + Inches(0.18), Inches(0.045), Inches(0.92), VERDE)

    # Descrição na direita
    add_text(s, desc, Inches(6.3), top + Inches(0.2), Inches(6.3), Inches(0.85),
             size=15, color=BRANCO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Gamificação
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "🎮  Sistema de Gamificação",
             "O diferencial do FIT TRACK — treinar vira um jogo com recompensas e progressão.")
numero_slide(s, 8)

gami = [
    ("⭐", "XP & Níveis",
     "• +50 XP por treino concluído\n• +50 XP ao completar missão\n• +20 XP meta de hidratação\n• +30 XP meta de passos\n• A cada 200 XP → sobe 1 nível"),
    ("📋", "Missões Diárias",
     "• Completar 1 treino do dia\n• Beber 2 litros de água\n• Caminhar 10.000 passos\nMissões se renovam todo dia"),
    ("🏆", "Conquistas (6 Badges)",
     "• Primeiro treino\n• 5 treinos realizados\n• 10 treinos realizados\n• 20 treinos realizados\n• Sequência de 3 dias\n• Sequência de 7 dias"),
    ("📈", "Estatísticas",
     "• Gráfico de atividade semanal\n• Contador de sequência (streak)\n• Total de treinos e calorias\n• Histórico completo de sessões"),
]

for i, (icon, title, desc) in enumerate(gami):
    col = i % 2
    row = i // 2
    lft = Inches(0.5 + col * 6.4)
    top = Inches(1.85 + row * 2.55)

    add_rect(s, lft, top, Inches(5.9), Inches(2.3), AZUL_CARD)
    add_rect(s, lft, top, Inches(0.14), Inches(2.3), VERDE)
    add_text(s, icon + "  " + title, lft + Inches(0.28), top + Inches(0.1),
             Inches(5.3), Inches(0.55), size=18, bold=True, color=VERDE)
    add_rect(s, lft + Inches(0.28), top + Inches(0.68), Inches(5.4), Inches(0.035), CINZA_TEXTO)
    add_text(s, desc, lft + Inches(0.28), top + Inches(0.76),
             Inches(5.4), Inches(1.45), size=13, color=BRANCO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Arquitetura
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "🏗️  Arquitetura da Aplicação",
             "Separação clara de responsabilidades: UI → Navegação → Estado → Serviços → Firebase")
numero_slide(s, 9)

layers = [
    (VERDE,                              "📱  INTERFACE  (React Native Screens)"),
    (AZUL_ESC,                           "🔀  NAVEGAÇÃO  (React Navigation — Stack + Tabs)"),
    (CINZA_MED,                          "⚙️  ESTADO GLOBAL  (Context API — AppContext)"),
    (RGBColor(0x0A, 0x29, 0x47),         "🔌  SERVIÇOS  (auth.js  |  firestore.js  |  storage.js)"),
    (PRETO,                              "☁️  FIREBASE  (Authentication  |  Firestore  |  Storage)"),
]

for i, (color, label) in enumerate(layers):
    top = Inches(1.5 + i * 1.05)
    add_rect(s, Inches(1.5), top, Inches(10.3), Inches(0.85), color)
    add_rect(s, Inches(1.5), top, Inches(0.12), Inches(0.85), VERDE)
    add_text(s, label, Inches(1.75), top + Inches(0.16), Inches(9.8), Inches(0.55),
             size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    if i < len(layers) - 1:
        add_text(s, "▼", Inches(6.55), top + Inches(0.85), Inches(0.7), Inches(0.22),
                 size=13, color=VERDE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Banco de Dados (Firestore)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)
slide_header(s, "🗄️  Banco de Dados — Firestore",
             "Estrutura NoSQL aninhada por usuário — dados isolados por UID, sem conflito entre contas.")
numero_slide(s, 10)

tree = [
    (0, "👤  users / {uid}", VERDE, 20, True),
    (1, "💪  workouts  →  id, name, exercises[], duration, calories", BRANCO, 15, False),
    (1, "📋  history   →  id, workoutName, sets, calories, date", BRANCO, 15, False),
    (1, "🎯  missions  →  id, title, xp, done, date", BRANCO, 15, False),
    (1, "👤  profile   →  name, email, weight, height, age, photoURL", BRANCO, 15, False),
]

box_top = Inches(1.9)
add_rect(s, Inches(0.8), box_top, Inches(11.7), Inches(4.0), AZUL_CARD)
add_rect(s, Inches(0.8), box_top, Inches(0.14), Inches(4.0), VERDE)

for i, (indent, text, color, size, bold) in enumerate(tree):
    add_text(s, text, Inches(1.1 + indent * 0.85), box_top + Inches(0.22 + i * 0.7),
             Inches(11.0), Inches(0.62), size=size, bold=bold, color=color)

perks = [
    "✅  Segurança: dados isolados por UID",
    "✅  Escalabilidade: NoSQL flexível, sem schema rígido",
    "✅  Tempo real: sincronização automática entre sessões",
    "✅  Offline: AsyncStorage garante persistência local",
]

for i, perk in enumerate(perks):
    col = i % 2
    row = i // 2
    add_text(s, perk,
             Inches(0.5 + col * 6.5), Inches(6.22 + row * 0.5),
             Inches(6.3), Inches(0.45),
             size=13, color=BRANCO)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusão
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)

add_rect(s, 0, 0, W, H, CINZA_MED)
add_rect(s, 0, 0, W, Inches(0.1), VERDE)
add_rect(s, 0, H - Inches(0.1), W, Inches(0.1), VERDE)
add_rect(s, 0, 0, Inches(0.18), H, VERDE)

numero_slide(s, 11)

add_text(s, "✅  Conclusão", Inches(0.8), Inches(0.5), Inches(12), Inches(0.8),
         size=40, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.8), Inches(1.32), Inches(5.73), Inches(0.055), VERDE)

points = [
    ("🎯", "Aplicativo funcional com 11 telas e fluxo completo de uso"),
    ("🔥", "Firebase como backend robusto — Auth + Firestore + Storage"),
    ("🎮", "Gamificação diferenciada: XP, níveis, missões e conquistas"),
    ("📐", "Arquitetura limpa com separação clara de responsabilidades"),
    ("📱", "Desenvolvido com React Native + Expo para Android"),
    ("👥", "Divisão organizada de tarefas entre os 4 integrantes"),
]

for i, (icon, point) in enumerate(points):
    top = Inches(1.65 + i * 0.78)
    add_rect(s, Inches(1.6), top, Inches(10.1), Inches(0.65), AZUL_CARD)
    add_rect(s, Inches(1.6), top, Inches(0.12), Inches(0.65), VERDE)
    add_text(s, icon, Inches(1.82), top + Inches(0.08), Inches(0.65), Inches(0.5), size=18)
    add_text(s, point, Inches(2.55), top + Inches(0.1), Inches(8.9), Inches(0.5),
             size=16, color=BRANCO)

add_rect(s, Inches(3.0), Inches(6.45), Inches(7.33), Inches(0.045), VERDE)
add_text(s, "Obrigado!  🏋️",
         Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.65),
         size=26, bold=True, color=VERDE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# Salvar
# ══════════════════════════════════════════════════════════════════
output = "c:/Users/canal/OneDrive/Desktop/Trabalho TI/FIT_TRACK_Apresentacao.pptx"
prs.save(output)
print(f"Arquivo salvo em: {output}")
